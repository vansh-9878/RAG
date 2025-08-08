import faiss,pickle
from concurrent.futures import ThreadPoolExecutor
from sentence_transformers import SentenceTransformer
from tqdm import tqdm
from langchain_core.tools import tool
from pptx import Presentation
# from localOCR import pdf_to_text
import numpy as np
import torch
import gc
import os
import threading
# import openpyxl
from docx import Document


# Set PyTorch CUDA memory allocation configuration to avoid fragmentation
# os.environ['PYTORCH_CUDA_ALLOC_CONF'] = 'expandable_segments:True'

# model = SentenceTransformer("BAAI/bge-small-en-v1.5")
# model = SentenceTransformer("all-MiniLM-L6-v2")
# Use GPU with fallback to CPU if CUDA is not available
try:
    model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")
except RuntimeError:
    model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2", device='cpu')
# model = SentenceTransformer("intfloat/e5-base-v2")
globalIndex=None
globalTexts=None
search_lock = threading.Lock()

arr=os.listdir('./vector')
arr=[item.split(".")[0] for item in arr]

def load_text_chunks(filepath, chunk_size=800,stride=200):
    with open(filepath, 'r', encoding='utf-8') as f:
        text = f.read()

    chunks = []

    for i in range(0, len(text), stride):
        chunk = text[i:i + chunk_size]
        if chunk:
            chunks.append(chunk)
    return chunks

def embed_in_batches(texts, model, batch_size=8, max_workers=2):
    embeddings = []

    def embed_batch(batch):
        # Clear GPU cache before processing batch
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        
        batch_embeddings = model.encode(batch, convert_to_numpy=True, normalize_embeddings=True)
        
        # Clear GPU cache after processing batch
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        
        return batch_embeddings

    # Process in smaller batches to avoid memory issues
    for i in tqdm(range(0, len(texts), batch_size), desc="Embedding"):
        batch = texts[i:i + batch_size]
        embeddings.append(embed_batch(batch))
        
        # Force garbage collection
        gc.collect()

    return np.vstack(embeddings)

def create_faiss_index(embeddings):
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)  
    index.add(embeddings)
    return index


def search_faiss(index, query, texts,top_k=25):
    with search_lock:  # Add thread safety
        query_embedding = model.encode([query], convert_to_numpy=True, normalize_embeddings=True)
        
        # Detailed debugging
        expected_dim = index.d
        actual_dim = query_embedding.shape[1]
        
        if expected_dim != actual_dim:
            return [("Error: Embedding dimension mismatch. Please regenerate the vector index.", 0.0)]
        
        try:
            D, I = index.search(query_embedding, top_k)
            results = [(texts[i], float(D[0][j])) for j, i in enumerate(I[0])]
            return results
        except Exception as e:
            import traceback
            traceback.print_exc()
            return [("Error during FAISS search", 0.0)]


def storeVectors(fileName):
    if fileName in arr:
        index = faiss.read_index(f"vector/{fileName}.faiss")
        with open(f"vector/{fileName}_texts.pkl", "rb") as f:
            texts = pickle.load(f)
    else:
        texts = load_text_chunks(fileName+".txt")

        embeddings = embed_in_batches(texts, model, batch_size=64, max_workers=8)
        index = create_faiss_index(embeddings)
        
        faiss.write_index(index, f"vector/{fileName}.faiss")

        with open(f"vector/{fileName}_texts.pkl", "wb") as f:
            pickle.dump(texts, f)
    
    return index,texts
    
@tool
def search(query,index,texts):
    """Search the document for the answer to the given query."""
    results = search_faiss(index, query, texts)
    return results


def ocrExcel():
    arr=[item for item in os.listdir('./unknownDoc') if item.endswith('.xlsx')]
    for file in arr:
        workbook = openpyxl.load_workbook(f"./unknownDoc/{file}")
        with open(f'{file[:-5]}.txt', 'w', encoding='utf-8') as txt_file:
            for sheet in workbook.worksheets:
                txt_file.write(f'--- Sheet: {sheet.title} ---\n')
                for row in sheet.iter_rows(values_only=True):
                    line = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    txt_file.write(line + '\n')
                txt_file.write('\n')
    
    
def ocrDocs():
    arr=[item for item in os.listdir('./unknownDoc') if item.endswith('.docx')]
    for file in arr:
        doc = Document(f"./unknownDoc/{file}")
        with open(f"{file[:-5]}.txt", 'w', encoding='utf-8') as txt_file:
            for para in doc.paragraphs:
                txt_file.write(para.text + '\n')
# ocrDocs()
def extract_text_from_pptx(pptx_path, output_txt_path):
    prs = Presentation(pptx_path)
    full_text = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = [f"--- Slide {slide_number} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        full_text.append("\n".join(slide_text))

    with open(output_txt_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(full_text))

    print(f"Text extracted and saved to '{output_txt_path}'")

extract_text_from_pptx("unknownDoc/Test Case HackRx.pptx","Test Case HackRx.txt")

# arr2 = [item for item in os.listdir('./') if item.endswith('.pdf') or item.endswith('.txt')]
# arr2 = [item[:-4] for item in arr2]

# for i in arr2:
#     storeVectors(i)